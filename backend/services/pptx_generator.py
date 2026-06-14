import io
import os
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

from schemas.pitch_deck import PitchDeck
from services.imagen_service import get_or_generate_slide_image

logger = logging.getLogger("pptx_generator")


def hex_to_rgb(hex_str: str, default_color: RGBColor) -> RGBColor:
    if not hex_str:
        return default_color
    hex_str = hex_str.strip().lstrip("#")
    try:
        if len(hex_str) == 6:
            r = int(hex_str[0:2], 16)
            g = int(hex_str[2:4], 16)
            b = int(hex_str[4:6], 16)
            return RGBColor(r, g, b)
    except Exception:
        pass
    return default_color


def get_category_theme(category: str) -> dict:
    """
    Get customized colors, backgrounds, and settings for the startup category.
    """
    cat = category.strip().lower()
    
    # Default values (SaaS Dark)
    bg = RGBColor(0x0B, 0x13, 0x2B)
    primary = RGBColor(0x48, 0xCA, 0xE4)
    secondary = RGBColor(0x00, 0xB4, 0xDB)
    text = RGBColor(0xFF, 0xFF, 0xFF)
    muted = RGBColor(0x9C, 0xA3, 0xAF)
    surface = RGBColor(0x1C, 0x25, 0x41)
    light_mode = False

    if "cyber" in cat:
        bg = RGBColor(0x05, 0x05, 0x05)
        primary = RGBColor(0x39, 0xFF, 0x14)  # Neon Green
        secondary = RGBColor(0x00, 0xF0, 0xFF)  # Neon Cyan
        surface = RGBColor(0x12, 0x22, 0x12)
    elif "fintech" in cat or "finance" in cat:
        bg = RGBColor(0xFF, 0xFF, 0xFF)
        primary = RGBColor(0x00, 0x35, 0x66)  # Deep Navy
        secondary = RGBColor(0xFF, 0xC3, 0x00)  # Gold
        text = RGBColor(0x12, 0x12, 0x12)
        muted = RGBColor(0x55, 0x55, 0x55)
        surface = RGBColor(0xF0, 0xF4, 0xF8)
        light_mode = True
    elif "agri" in cat or "farm" in cat or "crop" in cat:
        bg = RGBColor(0x0C, 0x1D, 0x0C)  # Dark Forest
        primary = RGBColor(0x9E, 0xF0, 0x1A)  # Lime
        secondary = RGBColor(0x70, 0xE0, 0x00)  # Green
        surface = RGBColor(0x18, 0x35, 0x18)
    elif "health" in cat or "med" in cat:
        bg = RGBColor(0xFF, 0xFF, 0xFF)
        primary = RGBColor(0x00, 0x80, 0x80)  # Teal
        secondary = RGBColor(0x20, 0xB2, 0xAA)  # Light Teal
        text = RGBColor(0x12, 0x12, 0x12)
        muted = RGBColor(0x55, 0x55, 0x55)
        surface = RGBColor(0xE6, 0xF2, 0xF2)
        light_mode = True
    elif "ai" in cat or "intel" in cat or "learning" in cat:
        bg = RGBColor(0x09, 0x07, 0x0F)  # Dark Purple Space
        primary = RGBColor(0xC7, 0x7D, 0xFF)  # Purple Accent
        secondary = RGBColor(0xE0, 0xAA, 0xFF)  # Secondary Soft Purple
        surface = RGBColor(0x1D, 0x14, 0x2D)
    elif "edu" in cat or "learn" in cat or "school" in cat:
        bg = RGBColor(0xFF, 0xFD, 0xF9)
        primary = RGBColor(0x80, 0x00, 0x20)  # Burgundy
        secondary = RGBColor(0xD4, 0xAF, 0x37)  # Warm Gold
        text = RGBColor(0x2B, 0x18, 0x10)
        muted = RGBColor(0x60, 0x50, 0x48)
        surface = RGBColor(0xF5, 0xEE, 0xE0)
        light_mode = True
    elif "logistics" in cat or "supply" in cat or "ship" in cat or "route" in cat:
        bg = RGBColor(0x1A, 0x1B, 0x22)
        primary = RGBColor(0xEF, 0x23, 0x3C)  # Red
        secondary = RGBColor(0x8D, 0x99, 0xAE)  # Steel Gray
        surface = RGBColor(0x2D, 0x31, 0x3D)
    elif "climate" in cat or "green" in cat or "carbon" in cat or "solar" in cat:
        bg = RGBColor(0x1E, 0x29, 0x3B)  # Slate
        primary = RGBColor(0x10, 0xB9, 0x81)  # Emerald
        secondary = RGBColor(0xA7, 0xF3, 0xD0)  # Sage
        surface = RGBColor(0x33, 0x41, 0x55)
    elif "commerce" in cat or "shop" in cat or "store" in cat:
        bg = RGBColor(0xFF, 0xFF, 0xFF)
        primary = RGBColor(0xFF, 0x4D, 0x4D)  # Coral
        secondary = RGBColor(0x4F, 0x46, 0xE5)  # Indigo
        text = RGBColor(0x1F, 0x29, 0x37)
        muted = RGBColor(0x6B, 0x72, 0x80)
        surface = RGBColor(0xF9, 0xFA, 0xFB)
        light_mode = True

    return {
        "bg": bg,
        "primary": primary,
        "secondary": secondary,
        "text": text,
        "muted": muted,
        "surface": surface,
        "light_mode": light_mode
    }


def style_cell(cell, text: str, bg_color: RGBColor, text_color: RGBColor, font_name: str, font_size: int = 11, bold: bool = False):
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg_color
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color


def build_pptx_presentation(deck_data: PitchDeck, state: dict, session_id: str, api_key: str | None) -> io.BytesIO:
    """
    Generate a highly visual, professional widescreen PPTX based on 
    the dynamic theme of the startup category.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Dynamic Theme Engine
    category = deck_data.category or "SaaS"
    theme = get_category_theme(category)
    
    bg_rgb = theme["bg"]
    primary_rgb = theme["primary"]
    secondary_rgb = theme["secondary"]
    text_rgb = theme["text"]
    muted_rgb = theme["muted"]
    surface_rgb = theme["surface"]
    light_mode = theme["light_mode"]
    
    font_name = "Segoe UI"
    blank_layout = prs.slide_layouts[6]

    for slide_data in deck_data.slides:
        slide = prs.slides.add_slide(blank_layout)
        
        # Background Style
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_rgb

        # Speaker notes
        if slide_data.speaker_notes:
            slide.notes_slide.notes_text_frame.text = slide_data.speaker_notes

        if slide_data.slide_number == 1:
            # ─── 1. Title/Cover Slide Layout ───
            # Large Title
            title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(1.8))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.title
            p.alignment = PP_ALIGN.CENTER
            p.font.name = font_name
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = primary_rgb
            
            # Elevator Subtitle
            subtitle_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.333), Inches(1.5))
            tf_sub = subtitle_box.text_frame
            tf_sub.word_wrap = True
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = slide_data.subtitle
            p_sub.alignment = PP_ALIGN.CENTER
            p_sub.font.name = font_name
            p_sub.font.size = Pt(20)
            p_sub.font.color.rgb = text_rgb
            
            # Bottom context notes
            context_box = slide.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(11.333), Inches(0.8))
            tf_ctx = context_box.text_frame
            tf_ctx.word_wrap = True
            p_ctx = tf_ctx.paragraphs[0]
            p_ctx.text = f"Startup Strategy Session   |   Category: {category}   |   Engine: exHacker"
            p_ctx.alignment = PP_ALIGN.CENTER
            p_ctx.font.name = font_name
            p_ctx.font.size = Pt(11)
            p_ctx.font.color.rgb = muted_rgb
            
        else:
            # ─── 2. Standard Split Layout ───
            # Slide Header Title
            header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9))
            tf_header = header_box.text_frame
            tf_header.word_wrap = True
            
            p_head = tf_header.paragraphs[0]
            p_head.text = slide_data.title
            p_head.font.name = font_name
            p_head.font.size = Pt(26)
            p_head.font.bold = True
            p_head.font.color.rgb = primary_rgb
            
            # Slide Subtitle
            p_sub = tf_header.add_paragraph()
            p_sub.text = slide_data.subtitle
            p_sub.font.name = font_name
            p_sub.font.size = Pt(13)
            p_sub.font.color.rgb = muted_rgb
            
            # Separation Line
            line_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(11.733), Inches(0.02)
            )
            line_shape.fill.solid()
            line_shape.fill.fore_color.rgb = secondary_rgb
            line_shape.line.fill.background()

            # Left Area Container Box (Top = 1.5", Width = 6.2", Height = 5.2")
            # Executive Summary (2-4 sentences)
            summary_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(6.2), Inches(0.9))
            tf_sum = summary_box.text_frame
            tf_sum.word_wrap = True
            p_sum = tf_sum.paragraphs[0]
            p_sum.text = slide_data.executive_summary
            p_sum.font.name = font_name
            p_sum.font.size = Pt(11)
            p_sum.font.italic = True
            p_sum.font.color.rgb = text_rgb
            p_sum.line_spacing = 1.25

            # Bullet points (6 to 10 points)
            bullet_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(6.2), Inches(3.2))
            tf_bullet = bullet_box.text_frame
            tf_bullet.word_wrap = True
            
            # Enforce 6-10 bullet points. Bullet text size is optimized to Pt(10) so it comfortably fits.
            for idx, bullet in enumerate(slide_data.bullets[:10]):
                p_bullet = tf_bullet.paragraphs[0] if idx == 0 else tf_bullet.add_paragraph()
                p_bullet.text = f"•   {bullet}"
                p_bullet.font.name = font_name
                p_bullet.font.size = Pt(10)
                p_bullet.font.color.rgb = text_rgb
                p_bullet.space_after = Pt(4)
                p_bullet.line_spacing = 1.15

            # Footer Highlight Container: Key Insight & Takeaway
            footer_card = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.8), Inches(6.2), Inches(1.1)
            )
            footer_card.fill.solid()
            footer_card.fill.fore_color.rgb = surface_rgb
            footer_card.line.color.rgb = primary_rgb
            footer_card.line.width = Pt(1.0)
            
            tf_f = footer_card.text_frame
            tf_f.word_wrap = True
            tf_f.margin_left = Inches(0.15)
            tf_f.margin_top = Inches(0.1)
            tf_f.margin_right = Inches(0.15)
            tf_f.margin_bottom = Inches(0.1)
            
            p_ins = tf_f.paragraphs[0]
            p_ins.text = f"KEY INSIGHT: {slide_data.key_insight}"
            p_ins.font.name = font_name
            p_ins.font.size = Pt(9.5)
            p_ins.font.bold = True
            p_ins.font.color.rgb = text_rgb
            p_ins.space_after = Pt(2)
            
            p_take = tf_f.add_paragraph()
            p_take.text = f"INVESTOR TAKEAWAY: {slide_data.investor_takeaway}"
            p_take.font.name = font_name
            p_take.font.size = Pt(9.5)
            p_take.font.italic = True
            p_take.font.color.rgb = secondary_rgb

            # Right Area Container Box (Left = 7.5", Top = 1.6", Width = 5.0", Height = 5.3")
            visual_left = Inches(7.5)
            visual_top = Inches(1.6)
            visual_width = Inches(5.0)
            visual_height = Inches(5.3)

            v_type = slide_data.visual_type.lower()
            
            if v_type == "image":
                # Imagen AI visual client
                image_path = get_or_generate_slide_image(
                    session_id=session_id,
                    slide_number=slide_data.slide_number,
                    prompt=slide_data.image_prompt,
                    api_key=api_key
                )
                
                if image_path and os.path.exists(image_path):
                    # 16:9 Image aligned centered in right column
                    img_height = Inches(2.81)
                    img_top = visual_top + (visual_height - img_height) / 2
                    slide.shapes.add_picture(
                        image_path, 
                        visual_left, 
                        img_top, 
                        width=visual_width, 
                        height=img_height
                    )
                else:
                    # Clean solid metric visual card fallback (never print prompts)
                    val_box = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE, visual_left, visual_top + Inches(0.6), visual_width, visual_height - Inches(1.2)
                    )
                    val_box.fill.solid()
                    val_box.fill.fore_color.rgb = surface_rgb
                    val_box.line.color.rgb = primary_rgb
                    val_box.line.width = Pt(1.5)
                    
                    tf_v = val_box.text_frame
                    tf_v.word_wrap = True
                    tf_v.margin_left = Inches(0.4)
                    tf_v.margin_top = Inches(0.5)
                    
                    p_v1 = tf_v.paragraphs[0]
                    p_v1.text = "STRATEGIC ANALYSIS"
                    p_v1.alignment = PP_ALIGN.CENTER
                    p_v1.font.name = font_name
                    p_v1.font.size = Pt(11)
                    p_v1.font.bold = True
                    p_v1.font.color.rgb = secondary_rgb
                    p_v1.space_after = Pt(20)
                    
                    p_v2 = tf_v.add_paragraph()
                    p_v2.text = "PROJECTION VALUE"
                    p_v2.alignment = PP_ALIGN.CENTER
                    p_v2.font.name = font_name
                    p_v2.font.size = Pt(12)
                    p_v2.font.color.rgb = muted_rgb
                    p_v2.space_after = Pt(10)

                    p_v3 = tf_v.add_paragraph()
                    p_v3.text = "10x Growth"
                    p_v3.alignment = PP_ALIGN.CENTER
                    p_v3.font.name = font_name
                    p_v3.font.size = Pt(38)
                    p_v3.font.bold = True
                    p_v3.font.color.rgb = text_rgb
                    
            elif v_type == "chart":
                # Market sizing column chart
                chart_data = CategoryChartData()
                chart_data.categories = ['SOM (Obtainable)', 'SAM (Serviceable)', 'TAM (Total Market)']
                chart_data.add_series('Sizing Projections ($B)', (1.5, 9.2, 54.0))
                
                chart_shape = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED, 
                    visual_left, 
                    visual_top + Inches(0.4), 
                    visual_width, 
                    visual_height - Inches(0.8), 
                    chart_data
                )
                chart = chart_shape.chart
                chart.has_legend = False
                
            elif v_type == "timeline":
                # Milestone timeline layout
                timeline_y = visual_top + visual_height / 2
                path_line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 
                    visual_left + Inches(0.2), 
                    timeline_y, 
                    visual_width - Inches(0.4), 
                    Inches(0.06)
                )
                path_line.fill.solid()
                path_line.fill.fore_color.rgb = secondary_rgb
                path_line.line.fill.background()
                
                milestones = [
                    ("Phase 1: Build", visual_left + Inches(0.3)),
                    ("Phase 2: Launch", visual_left + Inches(1.9)),
                    ("Phase 3: Scale", visual_left + Inches(3.5))
                ]
                
                for label, node_x in milestones:
                    node = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL, 
                        node_x, 
                        timeline_y - Inches(0.22), 
                        Inches(0.4), 
                        Inches(0.4)
                    )
                    node.fill.solid()
                    node.fill.fore_color.rgb = primary_rgb
                    node.line.color.rgb = text_rgb
                    node.line.width = Pt(1.5)
                    
                    label_box = slide.shapes.add_textbox(
                        node_x - Inches(0.55), 
                        timeline_y + Inches(0.35), 
                        Inches(1.5), 
                        Inches(1.5)
                    )
                    tf_label = label_box.text_frame
                    tf_label.word_wrap = True
                    p_lbl = tf_label.paragraphs[0]
                    p_lbl.text = label
                    p_lbl.alignment = PP_ALIGN.CENTER
                    p_lbl.font.name = font_name
                    p_lbl.font.size = Pt(11)
                    p_lbl.font.bold = True
                    p_lbl.font.color.rgb = text_rgb
                    
            elif v_type == "architecture":
                # System stack layer stack
                layers = [
                    ("CLIENT FRONTEND LAYER\n(Next.js App Client / Routing)", visual_top + Inches(0.2)),
                    ("CORE GATEWAY ROUTER\n(FastAPI / LangGraph Server)", visual_top + Inches(1.7)),
                    ("DATA STORAGE & MODELS\n(PostgreSQL / Supabase / AI Model)", visual_top + Inches(3.2))
                ]
                
                for title, y_pos in layers:
                    box = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE, 
                        visual_left + Inches(0.4), 
                        y_pos, 
                        visual_width - Inches(0.8), 
                        Inches(0.9)
                    )
                    box.fill.solid()
                    box.fill.fore_color.rgb = surface_rgb
                    box.line.color.rgb = primary_rgb
                    box.line.width = Pt(1.5)
                    
                    tf_box = box.text_frame
                    tf_box.word_wrap = True
                    p_box = tf_box.paragraphs[0]
                    p_box.text = title
                    p_box.alignment = PP_ALIGN.CENTER
                    p_box.font.name = font_name
                    p_box.font.size = Pt(9.5)
                    p_box.font.bold = True
                    p_box.font.color.rgb = text_rgb
                    
                    if y_pos != visual_top + Inches(3.2):
                        arrow = slide.shapes.add_shape(
                            MSO_SHAPE.DOWN_ARROW, 
                            visual_left + visual_width / 2 - Inches(0.2), 
                            y_pos + Inches(0.95), 
                            Inches(0.4), 
                            Inches(0.5)
                        )
                        arrow.fill.solid()
                        arrow.fill.fore_color.rgb = secondary_rgb
                        arrow.line.fill.background()
                        
            elif v_type == "comparison":
                # Competitive advantage table
                rows, cols = 4, 3
                grid_top = visual_top + Inches(0.4)
                grid_height = visual_height - Inches(0.8)
                
                table_shape = slide.shapes.add_table(
                    rows, 
                    cols, 
                    visual_left, 
                    grid_top, 
                    visual_width, 
                    grid_height
                )
                table = table_shape.table
                table.columns[0].width = Inches(1.8)
                table.columns[1].width = Inches(1.6)
                table.columns[2].width = Inches(1.6)
                
                headers = ["Key Moat / Advantage", "Competitors", "Us (Win)"]
                row_data = [
                    ["Automation Speed", "Manual / Delayed", "Instant Real-time"],
                    ["Feature Integration", "Fragmented / Silos", "Deep Unified Stack"],
                    ["Technical WoW", "Basic Forms", "AI Engine Node"]
                ]
                
                for col_idx, text_val in enumerate(headers):
                    style_cell(
                        table.cell(0, col_idx), 
                        text_val, 
                        bg_color=primary_rgb, 
                        text_color=text_rgb, 
                        font_name=font_name, 
                        font_size=9.5, 
                        bold=True
                    )
                    
                for row_idx, data in enumerate(row_data):
                    for col_idx, val in enumerate(data):
                        bg_color = surface_rgb
                        if col_idx == 2:
                            bg_color = RGBColor(0x1b, 0x3d, 0x1f) if not light_mode else RGBColor(0xdd, 0xf0, 0xde)
                        
                        style_cell(
                            table.cell(row_idx + 1, col_idx), 
                            val, 
                            bg_color=bg_color, 
                            text_color=text_rgb, 
                            font_name=font_name, 
                            font_size=9.5, 
                            bold=(col_idx == 2)
                        )
                        
            elif v_type == "funnel":
                # Go-To-Market Multi-stage funnel diagram
                funnel_stages = [
                    ("AWARENESS (Content GTM)", Inches(4.4), visual_left + Inches(0.3), visual_top + Inches(0.4)),
                    ("CONVERSION (POCTrial)", Inches(3.2), visual_left + Inches(0.9), visual_top + Inches(1.8)),
                    ("RETENTION (Enterprise Moat)", Inches(2.0), visual_left + Inches(1.5), visual_top + Inches(3.2))
                ]
                
                for stage_name, width, left_pos, y_pos in funnel_stages:
                    # Draw funnel deck trapezoid box using standard shapes
                    stage_box = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE, 
                        left_pos, 
                        y_pos, 
                        width, 
                        Inches(1.0)
                    )
                    stage_box.fill.solid()
                    stage_box.fill.fore_color.rgb = surface_rgb
                    stage_box.line.color.rgb = primary_rgb
                    stage_box.line.width = Pt(1.5)
                    
                    tf_stage = stage_box.text_frame
                    tf_stage.word_wrap = True
                    p_stage = tf_stage.paragraphs[0]
                    p_stage.text = stage_name
                    p_stage.alignment = PP_ALIGN.CENTER
                    p_stage.font.name = font_name
                    p_stage.font.size = Pt(9.5)
                    p_stage.font.bold = True
                    p_stage.font.color.rgb = text_rgb
                    
                    # Connection arrow downward
                    if y_pos != visual_top + Inches(3.2):
                        arrow = slide.shapes.add_shape(
                            MSO_SHAPE.DOWN_ARROW, 
                            visual_left + visual_width / 2 - Inches(0.15), 
                            y_pos + Inches(1.05), 
                            Inches(0.3), 
                            Inches(0.35)
                        )
                        arrow.fill.solid()
                        arrow.fill.fore_color.rgb = secondary_rgb
                        arrow.line.fill.background()
                        
            elif v_type == "diagram":
                # Business Model Revenue monetization diagram
                # Client Box (Left)
                left_box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, 
                    visual_left + Inches(0.2), 
                    visual_top + Inches(1.6), 
                    Inches(1.8), 
                    Inches(1.2)
                )
                left_box.fill.solid()
                left_box.fill.fore_color.rgb = surface_rgb
                left_box.line.color.rgb = primary_rgb
                left_box.line.width = Pt(1.5)
                
                tf_l = left_box.text_frame
                tf_l.word_wrap = True
                p_l = tf_l.paragraphs[0]
                p_l.text = "TARGET\nCUSTOMERS"
                p_l.alignment = PP_ALIGN.CENTER
                p_l.font.name = font_name
                p_l.font.size = Pt(10)
                p_l.font.bold = True
                p_l.font.color.rgb = text_rgb
                
                # Service Box (Right)
                right_box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, 
                    visual_left + Inches(3.0), 
                    visual_top + Inches(1.6), 
                    Inches(1.8), 
                    Inches(1.2)
                )
                right_box.fill.solid()
                right_box.fill.fore_color.rgb = surface_rgb
                right_box.line.color.rgb = secondary_rgb
                right_box.line.width = Pt(1.5)
                
                tf_r = right_box.text_frame
                tf_r.word_wrap = True
                p_r = tf_r.paragraphs[0]
                p_r.text = "OUR SAAS\nPLATFORM"
                p_r.alignment = PP_ALIGN.CENTER
                p_r.font.name = font_name
                p_r.font.size = Pt(10)
                p_r.font.bold = True
                p_r.font.color.rgb = text_rgb
                
                # Connecting flow arrows (Upper Right-ward and Lower Left-ward)
                arrow_top = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW, 
                    visual_left + Inches(2.1), 
                    visual_top + Inches(1.8), 
                    Inches(0.8), 
                    Inches(0.25)
                )
                arrow_top.fill.solid()
                arrow_top.fill.fore_color.rgb = primary_rgb
                arrow_top.line.fill.background()
                
                label_t = slide.shapes.add_textbox(
                    visual_left + Inches(2.0), 
                    visual_top + Inches(1.3), 
                    Inches(1.0), 
                    Inches(0.5)
                )
                label_t.text_frame.word_wrap = True
                p_lt = label_t.text_frame.paragraphs[0]
                p_lt.text = "Subscription"
                p_lt.font.name = font_name
                p_lt.font.size = Pt(8.5)
                p_lt.font.color.rgb = text_rgb
                
                arrow_bot = slide.shapes.add_shape(
                    MSO_SHAPE.LEFT_ARROW, 
                    visual_left + Inches(2.1), 
                    visual_top + Inches(2.35), 
                    Inches(0.8), 
                    Inches(0.25)
                )
                arrow_bot.fill.solid()
                arrow_bot.fill.fore_color.rgb = secondary_rgb
                arrow_bot.line.fill.background()
                
                label_b = slide.shapes.add_textbox(
                    visual_left + Inches(2.0), 
                    visual_top + Inches(2.6), 
                    Inches(1.0), 
                    Inches(0.5)
                )
                label_b.text_frame.word_wrap = True
                p_lb = label_b.text_frame.paragraphs[0]
                p_lb.text = "AI Value Out"
                p_lb.font.name = font_name
                p_lb.font.size = Pt(8.5)
                p_lb.font.color.rgb = text_rgb

            else:
                # Default blank backup card
                placeholder_card = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, visual_left, visual_top, visual_width, visual_height
                )
                placeholder_card.fill.solid()
                placeholder_card.fill.fore_color.rgb = surface_rgb
                placeholder_card.line.color.rgb = primary_rgb
                placeholder_card.line.width = Pt(1.5)

    # Save to binary output stream
    stream = io.BytesIO()
    prs.save(stream)
    stream.seek(0)
    return stream
